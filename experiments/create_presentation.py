from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT_DIR = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT_DIR / "wishlist_court_mcp_presentation.pptx"
VIDEO_PLACEHOLDER_SLIDE = 6

REFERENCE_FILES = [
    ROOT_DIR / "README.md",
    ROOT_DIR / "output" / "hermes_good_mcp_verdict.md",
    ROOT_DIR / "output" / "hermes_broken_mcp_verdict.md",
    ROOT_DIR / "output" / "orchestration_experiment_summary.md",
    ROOT_DIR / "experiments" / "presentation_slide_outline.md",
    ROOT_DIR / "HANDOFF_STATUS.md",
]

FONT_NAME = "Arial"
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(120, 120, 120)
LIGHT_GRAY = RGBColor(220, 220, 220)


def ensure_reference_files() -> None:
    missing = [str(path.relative_to(ROOT_DIR)) for path in REFERENCE_FILES if not path.exists()]
    if missing:
        raise FileNotFoundError("Missing reference files: " + ", ".join(missing))


def set_white_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def set_text_style(run, size: int, bold: bool = False, color=BLACK) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_notes(slide, lines: list[str]) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = "\n".join(lines)


def add_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.75), Inches(0.45), Inches(11.7), Inches(0.72))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    set_text_style(run, 33, bold=True)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.75), Inches(6.98), Inches(11.7), Inches(0.25))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    set_text_style(run, 10, color=GRAY)


def add_bullets(slide, bullets: list[str], top: float = 1.55, size: int = 24) -> None:
    box = slide.shapes.add_textbox(Inches(0.95), Inches(top), Inches(11.0), Inches(5.2))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True

    for idx, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(12)
        paragraph.font.name = FONT_NAME
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = BLACK


def add_small_label(slide, text: str, left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    set_text_style(run, 18, bold=True)


def add_plain_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    notes: list[str],
    footer: str = "",
    bullet_top: float = 1.55,
    bullet_size: int = 24,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide)
    add_title(slide, title)
    add_bullets(slide, bullets, top=bullet_top, size=bullet_size)
    if footer:
        add_footer(slide, footer)
    add_notes(slide, notes)


def set_cell_border(cell) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    for line_name in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tc_pr.find(qn(line_name))
        if existing is not None:
            tc_pr.remove(existing)
        line = OxmlElement(line_name)
        line.set("w", "9525")
        solid_fill = OxmlElement("a:solidFill")
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", "000000")
        solid_fill.append(srgb)
        line.append(solid_fill)
        tc_pr.append(line)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.05), Inches(11.5), Inches(0.9))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_run = title_frame.paragraphs[0].add_run()
    title_run.text = "Wishlist Court MCP"
    set_text_style(title_run, 46, bold=True)

    subtitle_box = slide.shapes.add_textbox(Inches(0.82), Inches(3.05), Inches(11.5), Inches(0.55))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    subtitle_run = subtitle_frame.paragraphs[0].add_run()
    subtitle_run.text = "Tool Quality가 LLM Reasoning에 미치는 영향 실험"
    set_text_style(subtitle_run, 24)

    meta_box = slide.shapes.add_textbox(Inches(0.84), Inches(4.65), Inches(8.0), Inches(0.95))
    meta_frame = meta_box.text_frame
    meta_frame.clear()
    for idx, line in enumerate(["생성모델응용", "이름: SimonWontae", "날짜: 2026-05-14"]):
        paragraph = meta_frame.paragraphs[0] if idx == 0 else meta_frame.add_paragraph()
        paragraph.text = line
        paragraph.space_after = Pt(5)
        paragraph.font.name = FONT_NAME
        paragraph.font.size = Pt(17)
        paragraph.font.color.rgb = BLACK

    add_notes(
        slide,
        [
            "이 프로젝트는 쇼핑 추천 앱이 아니라 MCP tool-use 실험임을 먼저 설명한다.",
            "같은 모델이라도 tool layer가 달라지면 reasoning이 어떻게 달라지는지 관찰했다고 소개한다.",
        ],
    )


def add_structure_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide)
    add_title(slide, "Wishlist Court MCP 구조")

    add_small_label(slide, "입력 데이터", 0.95, 1.45, 3.0, 0.35)
    add_bullets(slide, ["wishlist.csv", "owned_items.csv", "budget.json"], top=1.9, size=21)

    add_small_label(slide, "MCP 도구", 6.65, 1.45, 3.0, 0.35)
    tools = [
        "list_wishlist_items",
        "get_budget_status",
        "get_owned_items_by_category",
        "judge_purchase",
        "save_verdict_report",
    ]
    add_bullets(slide, tools, top=1.9, size=19)

    flow_box = slide.shapes.add_textbox(Inches(0.95), Inches(5.35), Inches(11.3), Inches(0.55))
    flow_frame = flow_box.text_frame
    flow_frame.clear()
    flow_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    flow_run = flow_frame.paragraphs[0].add_run()
    flow_run.text = "wishlist → budget 확인 → 보유품 조회 → 구매 판결 → markdown 저장"
    set_text_style(flow_run, 21, bold=True)

    bottom_box = slide.shapes.add_textbox(Inches(0.95), Inches(6.25), Inches(11.3), Inches(0.35))
    bottom_frame = bottom_box.text_frame
    bottom_frame.clear()
    bottom_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    bottom_run = bottom_frame.paragraphs[0].add_run()
    bottom_run.text = "LLM이 실제 상태(state)를 읽고 판단하도록 만드는 구조"
    set_text_style(bottom_run, 17)

    add_notes(
        slide,
        [
            "여기서 tool-use 흐름을 설명한다.",
            "단순 prompt engineering이 아니라 실제 데이터를 grounding한다고 설명한다.",
            "ReAct 스타일 흐름이라고 간단히 언급할 수 있다.",
        ],
    )


def add_demo_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide)
    add_title(slide, "실제 MCP Tool Call 데모")

    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.15),
        Inches(1.55),
        Inches(10.95),
        Inches(2.25),
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = WHITE
    placeholder.line.color.rgb = LIGHT_GRAY
    placeholder.line.width = Pt(1.25)
    text_frame = placeholder.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "[여기에 30초 데모 영상 삽입]"
    set_text_style(run, 28, bold=True)

    add_bullets(
        slide,
        [
            "1. MCP 등록 확인",
            "2. 정상 MCP 실행",
            "3. tool call 발생",
            "4. verdict 생성",
            "5. broken MCP 실행",
            "6. approve 편향 확인",
        ],
        top=4.15,
        size=20,
    )
    add_notes(
        slide,
        [
            "실제 Hermes tool call 장면이라고 설명한다.",
            "tool call 흐름을 짧게 언급한다.",
            "영상은 길게 설명하지 말고 빠르게 넘긴다.",
        ],
    )


def add_conclusion_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide)
    add_title(slide, "Orchestration 비교와 결론")

    add_bullets(
        slide,
        [
            "Single Agent: 가장 빠름, 토큰 효율 좋음, 기준 누락 위험",
            "Planner + Executor: 기준을 먼저 고정, 가장 안정적, 품질 최고",
            "Parallel sub-agent: 설명력 풍부, merge overhead 존재, 작은 작업에는 과함",
        ],
        top=1.35,
        size=21,
    )

    table = slide.shapes.add_table(4, 4, Inches(1.3), Inches(3.75), Inches(10.4), Inches(1.45)).table
    widths = [3.9, 2.0, 2.2, 1.6]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    headers = ["패턴", "Tokens", "Latency", "Quality"]
    rows = [
        ["Single", "4200", "18.2s", "4"],
        ["Planner+Executor", "6100", "27.5s", "5"],
        ["Parallel", "8900", "23.4s", "4"],
    ]
    for col, value in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = value
        cell.text_frame.paragraphs[0].font.bold = True
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            table.cell(row_idx, col_idx).text = value
    for row in table.rows:
        for cell in row.cells:
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.text_frame.paragraphs[0].font.name = FONT_NAME
            cell.text_frame.paragraphs[0].font.size = Pt(15)
            cell.text_frame.paragraphs[0].font.color.rgb = BLACK
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            set_cell_border(cell)

    takeaway_box = slide.shapes.add_textbox(Inches(0.95), Inches(6.05), Inches(11.4), Inches(0.55))
    takeaway_frame = takeaway_box.text_frame
    takeaway_frame.clear()
    paragraph = takeaway_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "LLM reasoning은 model 자체보다 tool layer의 quality와 structure에 크게 의존했다."
    set_text_style(run, 20, bold=True)

    add_notes(
        slide,
        [
            "orchestration은 보조 실험이라고 설명한다.",
            "Planner+Executor가 가장 안정적이었다고 결론 낸다.",
            "마지막 문장을 발표의 핵심 takeaway로 사용한다.",
        ],
    )


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_plain_slide(
        prs,
        "왜 Wishlist Court를 만들었는가?",
        [
            "LLM은 그럴듯한 소비 판단을 잘 생성한다.",
            "하지만 실제 예산, 기존 보유품, 구매 맥락을 모르면 쉽게 환각하거나 과소비 방향으로 흐를 수 있다.",
            "따라서 “위시리스트 + 예산 + 보유품”을 읽는 MCP 도구를 붙여 reasoning 변화를 관찰했다.",
            "예시 질문: “HHKB 키보드를 사도 될까?”",
        ],
        [
            "단순 챗봇이면 감성적 추천으로 흐를 수 있다고 설명한다.",
            "실제로는 예산과 기존 키보드 보유 여부가 중요하다고 설명한다.",
            "따라서 grounding용 MCP를 붙였다고 연결한다.",
        ],
        bullet_size=22,
    )
    add_structure_slide(prs)
    add_plain_slide(
        prs,
        "정상 MCP: Constraint-oriented reasoning",
        [
            "HHKB 키보드 → reject",
            "USB-C 허브 → approve",
            "에어팟 맥스 → cooldown",
            "남은 예산: 120,000 KRW",
            "기존 keyboard 보유품 2개 존재, 기존 USB-C 허브 상태 bad, urgency high 반영",
            "정상 MCP는 실제 예산과 보유품을 grounding하여 보수적 판단을 생성했다.",
        ],
        [
            "HHKB reject 사례를 중심으로 설명한다.",
            "기존 키보드가 이미 2개 있다는 점을 강조한다.",
            "USB-C 허브는 실제 문제 상황이라 approve된 점을 비교한다.",
            "constraint-oriented reasoning이라는 표현을 사용한다.",
        ],
        bullet_size=21,
    )
    add_plain_slide(
        prs,
        "Broken MCP: Reasoning distortion",
        [
            "실패 주입: fake budget, empty owned items, optimistic description",
            "결과: approve 5개, reject 0개",
            "정상 MCP: “예산 초과와 중복 보유 때문에 reject”",
            "Broken MCP: “사용 만족도를 높일 가능성이 있으므로 구매를 긍정적으로 검토”",
            "Broken MCP는 단순한 버그가 아니라 판단 기준 자체를 왜곡했다.",
        ],
        [
            "가장 중요한 슬라이드다.",
            "reasoning tone이 달라졌다는 점을 강조한다.",
            "constraint-oriented에서 desire-oriented reasoning으로 변한 것을 설명한다.",
            "tool description과 tool output이 모두 reasoning에 영향을 줬다고 설명한다.",
        ],
        bullet_size=21,
    )
    add_demo_slide(prs)
    add_conclusion_slide(prs)
    return prs


def main() -> None:
    ensure_reference_files()
    prs = build_presentation()
    prs.save(OUTPUT_PATH)
    print(f"Generated slides: {len(prs.slides)}")
    print(f"File path: {OUTPUT_PATH}")
    print(f"Video placeholder slide: {VIDEO_PLACEHOLDER_SLIDE}")


if __name__ == "__main__":
    main()
